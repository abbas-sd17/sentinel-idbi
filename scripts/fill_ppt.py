#!/usr/bin/env python3
"""Fill the mandatory IDBI Innovate PPT template with Sentinel content.

Template-compliance rules (host instruction: "You cannot change this
template ... everything should be same"):
- The template's mandated headings are NEVER modified or cleared. Content is
  added in NEW textboxes positioned below each heading.
- On the two slides whose text frame contains fill-in prompts (Team Details,
  Provide links), the heading paragraph is preserved and only the prompt
  lines are completed.
- Slide count and template imagery are untouched.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

TEMPLATE = Path("Prototype Submission Deck _ IDBI Innovate.pptx")
METRICS_PATH = Path("ml/reports/metrics.json")
SCREENSHOT_DIR = Path("docs/screenshots")
OUTPUT = Path("Sentinel_Submission_Deck.pptx")

BOTTOM_STRIP_IN = 5.52  # template's footer strip top edge
SLIDE_W_IN = 10.0


def _heading_shape(slide):
    """The template's heading textbox: first shape with non-empty text."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape
    return None


def _write_paragraphs(tf, text: str, font_size: int) -> None:
    """Write multi-line text into a text frame, one paragraph per line."""
    lines = text.split("\n")
    tf.text = lines[0]
    tf.paragraphs[0].font.size = Pt(font_size)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)


def _add_content_box(slide, text: str, font_size: int = 10) -> None:
    """Add a content textbox BELOW the template heading (heading untouched)."""
    heading = _heading_shape(slide)
    if heading is not None:
        top_in = heading.top / 914400 + heading.height / 914400 + 0.07
    else:
        top_in = 0.85
    height_in = max(BOTTOM_STRIP_IN - top_in - 0.05, 0.5)
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(top_in), Inches(SLIDE_W_IN - 0.7), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    _write_paragraphs(tf, text, font_size)


def _fill_prompts(slide, values: dict[str, str], font_size: int = 13) -> None:
    """Complete 'Prompt:'-style lines in the template's own text frame,
    preserving every paragraph that is not a known prompt (incl. heading)."""
    shape = _heading_shape(slide)
    if shape is None:
        return
    for p in shape.text_frame.paragraphs:
        raw = p.text.strip().rstrip(":").strip()
        if raw in values:
            filled = f"{p.text.strip().rstrip(':')}: {values[raw]}"
            for run in list(p.runs):
                run.text = ""
            if p.runs:
                p.runs[0].text = filled
            else:
                p.text = filled
            p.font.size = Pt(font_size)


def _segment_lines(metrics: dict) -> str:
    seg = metrics.get("segment_metrics") or []
    if not seg:
        return "  (per-segment metrics unavailable)"
    label = {
        "term_loan": "Term Loan", "cash_credit": "Cash Credit / WC",
        "lap": "LAP", "equipment_finance": "Equipment Finance",
    }
    return "\n".join(
        f"   {label.get(s['loan_type'], s['loan_type'])}: AUC {s['auc_roc']:.3f} / "
        f"KS {s['ks']:.3f} / Gini {s['gini']:.3f} (n={s['n']})"
        for s in seg
    )


def fill_deck(
    team_name: str,
    leader: str,
    github: str,
    demo_url: str,
    video_url: str,
) -> Path:
    prs = Presentation(str(TEMPLATE))
    metrics = {}
    if METRICS_PATH.exists():
        with open(METRICS_PATH, encoding="utf-8") as f:
            metrics = json.load(f)
    ew = metrics.get("early_warning", {})
    npv = metrics.get("npv", 0)
    npv_green = metrics.get("npv_green", 0)

    # Slide 1 (idx 0): fill the template's own prompt lines, heading preserved.
    _fill_prompts(
        prs.slides[0],
        {
            "Team name": team_name,
            "Team leader name": leader,
            "Problem Statement": "Track 04 - Default Prediction Model",
        },
    )

    # Slide 14 (idx 13): fill the links prompts in place, heading preserved.
    _fill_prompts(
        prs.slides[12],
        {
            "GitHub Public Repository": github,
            "Demo Video Link (3 Minutes)": video_url,
            "Final Product Link": demo_url,
        },
        font_size=12,
    )

    content = {
        1: (
            "Sentinel is an explainable MSME loan default early-warning system that "
            "predicts the probability of default 12 months in advance. It fuses the three "
            "input classes the problem statement mandates - borrower behavior, bank-internal "
            "data, and public-domain data (electricity consumption, EPFO payroll, Udyam "
            "status) - plus unstructured RM notes, into one calibrated scoring engine. "
            "Loan officers get a unified PD, a RAG (Red/Amber/Green) bucket, an RBI SMA "
            "early-stress view, an internal rating grade, evidence-gated IFRS-9 staging, "
            "SHAP reason codes, and a decision panel that records their disposition to an "
            "audit trail - proactive intervention months before an account slips to NPA."
        ),
        2: (
            "Opportunities:\n"
            "- Incumbent model: 16-22% accuracy, structured-only, ~3-month warning; Sentinel warns 12 months out\n"
            "- One engine + per-segment thresholds ends the fragmented methodology across loan types\n"
            "- RBI-ready governance: SMA/EWS alignment, audit trail, drift monitoring, model card\n"
            "How it solves the problem:\n"
            "- Calibrated gradient boosting on structured + behavioral + public-domain + NLP signals\n"
            "- Delivering the >90% intent: 95.6% of accounts the model clears stay good over the next\n"
            "  12 months (NPV); Red-tier precision 40.7% is ~2x the incumbent's 16-22% accuracy\n"
            "- Human-in-the-loop by construction: AI advises with reasons; the officer decides, and\n"
            "  every decision is recorded (RBI expectations)\n"
            "USP:\n"
            "- Honest, leakage-free held-out validation (AUC 0.786 / KS 0.43 - the band real bank PD\n"
            "  models occupy), with threshold selected on a separate validation split\n"
            "- Live-regime outputs: SMA-0/1/2 view + evidence-gated Ind-AS 109 readiness + PD x LGD x EAD ECL\n"
            "- Standardized reason-code taxonomy (BEH/FIN/BUR/CMP/TXT/PUB) mapping to RBI EWS signals\n"
            "- 100% open source, Docker/AWS-portable, synthetic DPDP-safe data ready for sandbox swap"
        ),
        3: (
            "- 12-month calibrated PD (0-100%) per account, single-account and validated batch CSV scoring\n"
            "- RAG (Red/Amber/Green) bucketing with per-loan-type thresholds + G1-G10 rating masterscale\n"
            "- RBI SMA early-stress view (Standard/SMA-0/1/2, trailing-12m worst) - the live IRAC regime\n"
            "- Evidence-gated IFRS-9/Ind-AS 109 staging (Stage 3 only on 90+ DPD) + ECL = PD x LGD x EAD\n"
            "- SHAP reason codes with bank-auditable taxonomy incl. public-domain signals (PUB-01/02/03)\n"
            "- Officer decision capture (acknowledge/override/escalate + note) logged to the audit trail\n"
            "- 12-month hazard-curve decomposition consistent with the headline PD\n"
            "- Portfolio dashboard: KPIs, RAG/SMA breakdowns, sector heatmap, paginated watchlist\n"
            "- Batch upload validates every row and reports skipped rows + neutrally-imputed columns\n"
            "- PSI drift monitoring endpoint + JSONL audit trail (model governance)"
        ),
        4: (
            "1. Observation snapshot of each MSME loan account (Stage 2: monthly panel from sandbox)\n"
            "2. Feature engineering: structured + behavioral + public-domain (electricity/EPFO/Udyam) + NLP\n"
            "3. Segment assignment (loan type x borrower profile x enterprise size)\n"
            "4. Calibrated model scores PD of default within 12 months\n"
            "5. Per-segment thresholds map PD to RAG; DPD evidence gates SMA category and IFRS-9 stage\n"
            "6. SHAP produces top reason codes with direction (increases/decreases risk)\n"
            "7. Loan officer reviews, then acknowledges / overrides / escalates - decision recorded via\n"
            "   POST /decisions to the same audit trail as the score (human-in-the-loop, implemented)\n"
            "8. Portfolio dashboard aggregates risk; PSI monitoring watches population drift"
        ),
        5: (
            "Live prototype (see Final Product Link):\n"
            "- Portfolio Dashboard: KPI cards, RAG pie, RBI SMA strip, sector heatmap, paginated account table\n"
            "- Account Detail: PD gauge, RAG badge, SMA + IFRS-9 chips, hazard curve, SHAP reason codes,\n"
            "  officer decision panel (acknowledge/override/escalate + note)\n"
            "- Batch Upload: CSV drag-and-drop with per-row validation, skipped-row report and\n"
            "  imputed-column warnings"
        ),
        6: (
            "Architecture (3-tier, open-source, AWS-portable):\n"
            "Frontend: Next.js 14 + Tailwind (Netlify / Vercel / AWS Amplify)\n"
            "API: FastAPI + Docker (Hugging Face Spaces / self-hosted / AWS ECS Fargate)\n"
            "ML: scikit-learn HistGradientBoosting + isotonic calibration; TF-IDF NLP with fit-time-frozen\n"
            "     normalization (identical single-account and batch scores)\n"
            "Data: synthetic MSME parquet (swaps to IDBI sandbox APIs in Stage 2)\n"
            "Explainability: SHAP with disclosed fallback; method recorded per decision\n"
            "Governance: JSONL audit trail (scores + officer decisions), PSI drift endpoint\n"
            "Stage 2 AWS mapping: API -> ECS Fargate | artifacts -> S3 | data -> sandbox APIs |\n"
            "frontend -> Amplify | entry -> API Gateway + ALB"
        ),
        7: (
            "Technologies (all open source):\n"
            "- Python 3.12, scikit-learn, pandas, numpy, SHAP\n"
            "- FastAPI, uvicorn, pydantic (schema-validated scoring, incl. batch)\n"
            "- Next.js 14, React 18, Tailwind CSS, Recharts\n"
            "- Docker, docker-compose; GitHub Actions CI\n"
            "- Deployment: Hugging Face Spaces (API) + Netlify (frontend); fully self-hostable on-prem\n"
            "- Stage 2 target: AWS (ECS, S3, API Gateway, Amplify) - matches the bank's cloud posture\n"
            "- Compliance: DPDP-safe synthetic data, RBI SMA/EWS alignment, model card, audit trail"
        ),
        8: (
            "Estimated AWS pilot cost (monthly):\n"
            "- ECS Fargate (1 task): ~$15-25\n"
            "- S3 storage: ~$1-5\n"
            "- API Gateway: ~$3-10\n"
            "- Amplify hosting: ~$0-15\n"
            "- Total pilot: ~$20-55/month\n"
            "- Stage 1 prototype: free/open-source (Hugging Face Spaces + Netlify), or on-prem docker-compose"
        ),
        9: None,  # screenshots added as images below
        10: (
            f"Held-out test set, n={metrics.get('n_samples', 'N/A')} (synthetic; validates the pipeline - Stage 2 revalidates on IDBI data):\n"
            f"AUC {metrics.get('auc_roc', 'N/A')} | KS {metrics.get('ks_statistic', 'N/A')} | Gini {metrics.get('gini', 'N/A')} | "
            f"PR-AUC {metrics.get('pr_auc', 'N/A')} (base ~{round(metrics.get('default_rate', 0) * 100, 1)}%) | Brier {metrics.get('brier_score', 'N/A')} | "
            f"top-decile lift {metrics.get('lift_at_10pct', 'N/A')}x\n"
            f"Delivering the >90% intent: NPV {npv:.1%} of cleared accounts stay good; green-clearance "
            f"{npv_green:.1%};\n"
            f"Red-tier precision {metrics.get('precision_at_red', 'N/A')} = ~2x the incumbent's 16-22%; recall (defaulters) {metrics.get('recall', 'N/A')}\n"
            f"Early warning: {ew.get('defaulters_flagged_pct', 'N/A')}% of eventual defaulters on the watchlist, avg "
            f"{ew.get('avg_lead_time_months', 'N/A')} months ahead (lead time illustrative - synthetic timing)\n"
            f"Consistent across products (held-out):\n{_segment_lines(metrics)}\n"
            f"Method notes: threshold selected on a separate validation split (test never used for selection);\n"
            f"in-sample reads ~0.87 AUC vs held-out {metrics.get('auc_roc', 'N/A')} - we report held-out; raw accuracy "
            f"{metrics.get('raw_accuracy', 'N/A')} disclosed\n"
            f"(uninformative alone at a 9% default rate). Charts below: ROC / KS separation / calibration."
        ),
        11: (
            "Future development (Stage 2+):\n"
            "- IDBI sandbox migration: map columns, retrain, out-of-time + vintage validation on real data\n"
            "- Live public-domain feeds: state discom electricity, EPFO API, Udyam portal, MCA/NCLT watch\n"
            "- Collateral-adjusted product-level LGDs for ECL; champion-challenger + periodic recalibration\n"
            "- CBS loan-master integration; monthly observation panel for true lead-time measurement\n"
            "- RBI EWS/RFA signal-library mapping expansion; maker-checker workflow on overrides\n"
            "- Domain-tuned NLP embeddings for RM notes; retail/corporate segment expansion\n"
            "- Real-time streaming scoring (Kafka/Kinesis); pilot deployment in IDBI ecosystem"
        ),
    }

    for idx, text in content.items():
        if text is None or idx >= len(prs.slides):
            continue
        size = 10 if idx not in (1, 5, 8) else 11
        if idx == 10:
            size = 9
        _add_content_box(prs.slides[idx], text, font_size=size)

    # Slide 10 (idx 9): prototype screenshots in a grid, if captured.
    shots = sorted(SCREENSHOT_DIR.glob("*.png")) if SCREENSHOT_DIR.exists() else []
    if shots and len(prs.slides) > 9:
        slide = prs.slides[9]
        max_h = Inches(1.85)
        positions = [
            (0.35, 1.55), (5.05, 1.55),
            (0.35, 3.55), (5.05, 3.55),
        ]
        for (x, y), img in zip(positions, shots[:4]):
            slide.shapes.add_picture(str(img), Inches(x), Inches(y), height=max_h)
    elif len(prs.slides) > 9:
        _add_content_box(
            prs.slides[9],
            "Screenshots: portfolio dashboard / account detail with decision panel / "
            "batch upload validation. (Capture from the running prototype into "
            "docs/screenshots/*.png and re-run this script to embed.)",
            font_size=11,
        )

    # Slide 11 (idx 10): benchmark charts sized to FIT the slide (no overflow).
    reports = Path("ml/reports")
    if len(prs.slides) > 10:
        bench = prs.slides[10]
        chart_h = Inches(1.35)
        for i, img_name in enumerate(["roc_curve.png", "ks_curve.png", "calibration_curve.png"]):
            img_path = reports / img_name
            if img_path.exists():
                bench.shapes.add_picture(
                    str(img_path), Inches(0.55 + i * 3.2), Inches(4.05), height=chart_h
                )

    prs.save(str(OUTPUT))
    print(f"Saved filled deck to {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--team-name", default="ax3nr1x")
    parser.add_argument("--leader", default="Syed Abbas Raza")
    parser.add_argument("--github", default="https://github.com/abbas-sd17/sentinel-idbi")
    parser.add_argument("--demo-url", default="https://sentinel-idbi.netlify.app")
    parser.add_argument("--video-url", default="https://youtu.be/DLrS6FQ7Jcs")
    args = parser.parse_args()
    fill_deck(args.team_name, args.leader, args.github, args.demo_url, args.video_url)
